"""
Builds the 15-slide main deck as a .pptx.

Copy comes from DECK_STORYBOARD.md and is reproduced as written. Numbers are
not calculated here, they are transcribed from the storyboard, which in turn
traces to cost_model.py. If a number needs changing, change it in the engine
and the storyboard first, then here.

Run from the recommendation/ directory:
    ../.venv/bin/python build_deck.py

Writes ../deck/NovaCorp_deck.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# --------------------------------------------------------------------------
# SET THIS BEFORE THE FINAL EXPORT
# --------------------------------------------------------------------------
TEAM_NAME = "O for 4"
DATE_LINE = "15 August 2026"

# --------------------------------------------------------------------------
HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
OUT_DIR = ROOT / "deck"
OUT_DIR.mkdir(exist_ok=True)

FIG = ROOT / "figures"
EFIG = ROOT / "ethics_finance" / "figures_ethics"
P2FIG = ROOT / "explore_pt2" / "part2b_outputs" / "figures"
DECKFIG = OUT_DIR / "figures"

PURPLE = RGBColor(0xA1, 0x00, 0xFF)
DEEP   = RGBColor(0x46, 0x00, 0x73)
TEAL   = RGBColor(0x00, 0xB7, 0xC3)
CORAL  = RGBColor(0xFF, 0x6B, 0x6B)
GOLD   = RGBColor(0xFF, 0xB3, 0x00)
INK    = RGBColor(0x22, 0x22, 0x2A)
GREY   = RGBColor(0x5A, 0x5A, 0x66)
LIGHT  = RGBColor(0xEC, 0xEC, 0xF1)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.62)
BODY_W = W - 2 * MARGIN

FONT = "Calibri"


# --------------------------------------------------------------------------
# helpers
# --------------------------------------------------------------------------
def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, color=INK, bold=False, italic=False, space_after=6,
         align=PP_ALIGN.LEFT, first=False, line_spacing=1.0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    _runs(p, text, size, color, bold, italic)
    return p


def _runs(p, text, size, color, bold, italic):
    """Splits on ** ** so a bullet can carry bold emphasis inline."""
    for i, chunk in enumerate(text.split("**")):
        if not chunk:
            continue
        r = p.add_run()
        r.text = chunk
        r.font.size = Pt(size)
        r.font.name = FONT
        r.font.italic = italic
        r.font.bold = bold or (i % 2 == 1)
        r.font.color.rgb = DEEP if (i % 2 == 1 and not bold) else color


def picture(slide, path, box_x, box_y, box_w, box_h):
    """Fits an image inside a box, preserving aspect, centred."""
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(box_w / iw, box_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    x = int(box_x + (box_w - w) / 2)
    y = int(box_y + (box_h - h) / 2)
    return slide.shapes.add_picture(str(path), x, y, w, h)


def chrome(slide, eyebrow, headline, number, footnote=None, hl_size=27):
    """Standard slide furniture: eyebrow, headline, rule, page number, source."""
    rect(slide, 0, 0, W, Inches(0.075), fill=PURPLE)

    tf = textbox(slide, MARGIN, Inches(0.40), BODY_W, Inches(0.24))
    para(tf, eyebrow.upper(), 10.5, GREY, bold=True, first=True, space_after=0)

    tf = textbox(slide, MARGIN, Inches(0.70), BODY_W, Inches(0.95))
    para(tf, headline, hl_size, DEEP, bold=True, first=True, space_after=0,
         line_spacing=1.06)

    if footnote:
        tf = textbox(slide, MARGIN, Inches(6.92), Inches(11.2), Inches(0.42))
        para(tf, footnote, 8, GREY, italic=True, first=True, space_after=0,
             line_spacing=1.12)

    tf = textbox(slide, W - MARGIN - Inches(0.6), Inches(6.95), Inches(0.6), Inches(0.25))
    para(tf, str(number), 10, GREY, first=True, space_after=0, align=PP_ALIGN.RIGHT)


def bullets(slide, x, y, w, h, items, size=13.5, gap=11):
    tf = textbox(slide, x, y, w, h)
    for i, t in enumerate(items):
        p = para(tf, t, size, INK, first=(i == 0), space_after=gap,
                 line_spacing=1.18)
        p.level = 0
    return tf


def table(slide, x, y, w, rows, col_w, header_fill=DEEP, row_h=Inches(0.38),
          size=12, head_size=11.5, emphasis_rows=()):
    n_r, n_c = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_r, n_c, x, y, w, row_h * n_r)
    tbl = shape.table
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = cw
    for i, row in enumerate(rows):
        tbl.rows[i].height = row_h
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = ""
            c.margin_left = Inches(0.10)
            c.margin_right = Inches(0.08)
            c.margin_top = c.margin_bottom = Inches(0.03)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(val)
            r.font.name = FONT
            r.font.size = Pt(head_size if i == 0 else size)
            r.font.bold = (i == 0) or (i in emphasis_rows)
            r.font.color.rgb = WHITE if i == 0 else INK
            c.fill.solid()
            if i == 0:
                c.fill.fore_color.rgb = header_fill
            elif i in emphasis_rows:
                c.fill.fore_color.rgb = RGBColor(0xF4, 0xE9, 0xFF)
            else:
                c.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(0xF8, 0xF8, 0xFB)
    return tbl


def callout(slide, x, y, w, h, text, accent=PURPLE, size=12.5):
    rect(slide, x, y, Inches(0.055), h, fill=accent)
    tf = textbox(slide, x + Inches(0.22), y, w - Inches(0.22), h,
                 anchor=MSO_ANCHOR.MIDDLE)
    para(tf, text, size, INK, italic=True, first=True, space_after=0,
         line_spacing=1.16)


# ==========================================================================
# SLIDES
# ==========================================================================
def slide_01(prs):
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=DEEP)
    rect(s, 0, Inches(3.02), Inches(1.5), Inches(0.07), fill=PURPLE)

    tf = textbox(s, MARGIN, Inches(2.05), Inches(11.4), Inches(0.5))
    para(tf, "ACCENTURE × SUBAA  ·  NOVACORP PEOPLE ANALYTICS CHALLENGE",
         12, RGBColor(0xC9, 0xA8, 0xEA), bold=True, first=True, space_after=0)

    tf = textbox(s, MARGIN, Inches(3.35), Inches(11.6), Inches(1.5))
    para(tf, "What's driving the $42M", 46, WHITE, bold=True, first=True,
         space_after=2, line_spacing=1.02)
    para(tf, "and where to act first", 46, PURPLE, bold=True, space_after=0,
         line_spacing=1.02)

    tf = textbox(s, MARGIN, Inches(5.55), Inches(11.4), Inches(0.4))
    para(tf, f"{TEAM_NAME}  ·  {DATE_LINE}", 13,
         RGBColor(0xC9, 0xA8, 0xEA), first=True, space_after=0)


def slide_02(prs):
    s = blank(prs)
    chrome(s, "The question you asked us",
           "You asked what's driving $42M.\nWe found you're counting $14M of it.", 2)

    bullets(s, MARGIN, Inches(2.30), Inches(7.1), Inches(3.4), [
        "Finance sized the annual people cost at **~$42M** across three buckets: "
        "regrettable attrition ($22–25M), disengagement ($12–15M), hiring "
        "inefficiency ($4–6M).",
        "We did not set out to re-litigate that number. We set out to find which "
        "parts are **tractable**.",
        "What we found first was a **measurement problem**, and it changes the "
        "shape of everything underneath it.",
    ], size=14.5, gap=15)

    x = Inches(8.15)
    rect(s, x, Inches(2.30), Inches(4.55), Inches(3.35), fill=RGBColor(0xF7, 0xF3, 0xFC))
    tf = textbox(s, x + Inches(0.4), Inches(2.62), Inches(3.8), Inches(0.4))
    para(tf, "WHERE WE'RE GOING", 10, GREY, bold=True, first=True, space_after=14)
    for i, (n, t) in enumerate([
            ("1", "What your data is hiding"),
            ("2", "Where the loss actually sits"),
            ("3", "What it's worth if you act")]):
        yy = Inches(3.20 + i * 0.70)
        tf = textbox(s, x + Inches(0.4), yy, Inches(0.4), Inches(0.4))
        para(tf, n, 20, PURPLE, bold=True, first=True, space_after=0)
        tf = textbox(s, x + Inches(0.85), yy + Inches(0.06), Inches(3.4), Inches(0.4))
        para(tf, t, 14, INK, first=True, space_after=0)

    callout(s, MARGIN, Inches(6.05), Inches(12.0), Inches(0.62),
            "All of it uses your own numbers and Finance's own formulas.")


def slide_03(prs):
    s = blank(prs)
    chrome(s, "The headline",
           "You count $14.3M of regrettable attrition. The real figure is $45.1M.", 3,
           "Source: attrition_log.csv, brief §6 constants. Rates are annualised over the "
           "two-year window. Defensible range $34–56M, see appendix A5.")
    bullets(s, MARGIN, Inches(2.05), Inches(4.35), Inches(4.4), [
        "HR's regrettable_flag fires on **153 departures** over two years, "
        "which is **$14.3M/yr**.",
        "Apply a value-based definition, every High Performer or Outstanding "
        "employee who **chose** to leave, and it is **499 departures → $45.1M/yr**.",
        "**371 high-performing people** walked out and were never recorded as a loss.",
    ], size=13.5, gap=14)
    picture(s, EFIG / "E1_regrettable_gap.png",
            Inches(5.35), Inches(1.95), Inches(7.35), Inches(4.75))
    callout(s, MARGIN, Inches(6.05), Inches(4.35), Inches(0.72),
            "The same replacement-cost formula Finance gave us, applied to a "
            "population their flag excludes.", accent=CORAL, size=11)


def slide_04(prs):
    s = blank(prs)
    chrome(s, "Why the flag misses them",
           "The flag is a synonym for “Outstanding”, not a measure of value lost.", 4,
           "Source: attrition_log.csv, regrettable_flag × performance_band_at_exit × pathway. "
           "n = 1,400 exits over two years.")
    bullets(s, MARGIN, Inches(2.05), Inches(4.35), Inches(4.4), [
        "Flag rate by rating: **Outstanding 61–72%**. High Performer **15% (pull), "
        "0% (push)**. Meets Expectations 12% / 1%.",
        "Of the **168 High Performers who were managed out, zero** were flagged "
        "regrettable. Not a low rate. Zero.",
        "The flag is applied **after** the exit, by the same function that approved "
        "it. It records whether a departure was **convenient**, not whether it was "
        "**costly**.",
    ], size=13, gap=13)
    picture(s, EFIG / "E2_flag_never_fires.png",
            Inches(5.35), Inches(1.95), Inches(7.35), Inches(4.75))
    callout(s, MARGIN, Inches(6.02), Inches(4.35), Inches(0.80),
            "Not bad faith. A retrospective, self-assessed metric that structurally "
            "exonerates the decision-maker is a justification, not a measurement.",
            accent=CORAL, size=10.5)


def slide_05(prs):
    s = blank(prs)
    chrome(s, "The $42M, restated",
           "It isn't $42M. On your own formulas it's $65M, and the shape is different too.",
           5, "All figures from cost_model.py on the brief's §6 constants. "
              "Same formula as Finance, different population.", hl_size=25)
    picture(s, EFIG / "E6_buckets_restated.png",
            Inches(0.55), Inches(1.85), Inches(7.35), Inches(4.35))
    rows = [
        ["Bucket", "Finance", "On the data"],
        ["Regrettable attrition", "$22–25M", "$45.1M"],
        ["Disengagement", "$12–15M", "$17.8M"],
        ["Hiring inefficiency", "$4–6M", "$2.3M"],
        ["Total", "$38–46M", "$65.2M"],
    ]
    table(s, Inches(8.15), Inches(2.10), Inches(4.6), rows,
          [Inches(2.15), Inches(1.15), Inches(1.30)], emphasis_rows=(4,))
    callout(s, Inches(8.15), Inches(4.60), Inches(4.6), Inches(1.15),
            "Only one of these moves in your favour, and it's the smallest one. "
            "The two that moved against you were both sized with a definition that "
            "excluded the expensive cases.", size=11)


def slide_06(prs):
    s = blank(prs)
    chrome(s, "Where the loss concentrates",
           "Entity_B loses people at twice Entity_A's rate, and it's the one still "
           "on its own HR system.", 6,
           "Rates are annualised voluntary exits ÷ active headcount (n = 12,003), the same "
           "convention used on every rate in this deck. Chi-square, Entity_B vs Entity_A, "
           "p = 8×10⁻⁹.", hl_size=25)
    bullets(s, MARGIN, Inches(2.15), Inches(4.35), Inches(4.2), [
        "NovaCorp-Origin **4.5%** · Entity_A **3.6%** · **Entity_B 7.0%** · "
        "Entity_C **4.6%**, per year.",
        "Entity_B vs Entity_A on 1,601 and 1,804 active staff. The most "
        "statistically solid finding in the analysis.",
        "Entity_B was acquired in FY2023 and **still runs on BambooHR**, two years "
        "on. Entity_A was absorbed onto the core system and is now the healthiest "
        "cohort in the company.",
    ], size=13, gap=13)
    picture(s, FIG / "03_attrition_by_legacy_entity.png",
            Inches(5.35), Inches(1.95), Inches(7.35), Inches(4.60))
    callout(s, MARGIN, Inches(5.90), Inches(4.35), Inches(0.90),
            "Entity_B isn't where the most dollars sit. NovaCorp-Origin holds $29.2M "
            "of the $45.1M on headcount alone. Entity_B is where the leverage is: "
            "$8.4M in 1,601 people you already know how to fix.", size=10.5)


def slide_07(prs):
    s = blank(prs)
    chrome(s, "It isn't pay, and it isn't their managers",
           "Six of eight engagement dimensions are fine. Two collapsed.", 7,
           "Employee-level means across five waves, Welch's t-test. Entity_A n = 1,938, "
           "Entity_B n = 1,697 responders. Composite index: Entity_B 3.280 vs Entity_A 3.366.")
    bullets(s, MARGIN, Inches(2.05), Inches(4.35), Inches(4.4), [
        "**Manager effectiveness**: 3.344 vs 3.371, **no difference** (p = 0.35).",
        "**Psychological safety**: 3.351 vs 3.354, **no difference** (p = 0.91).",
        "**Senior leadership trust −0.301** (p = 1.5×10⁻²⁵). "
        "**Purpose & meaning −0.300** (p = 2.6×10⁻²⁵).",
        "And it isn't money. Entity_B sits at **0.961** compa-ratio against "
        "NovaCorp-Origin's **0.937**. They are paid better and leave anyway.",
    ], size=12.5, gap=11)
    picture(s, P2FIG / "13_entity_b_dimension_diagnosis.png",
            Inches(5.35), Inches(1.90), Inches(7.35), Inches(4.55))
    callout(s, MARGIN, Inches(5.72), Inches(4.35), Inches(1.05),
            "Why it never showed on your dashboard: the composite index averages all "
            "eight dimensions, so two collapsed scores are diluted by six healthy ones. "
            "A gap of 0.085 that looks like nothing.", accent=GOLD, size=10.5)


def slide_08(prs):
    s = blank(prs)
    chrome(s, "They arrived this way, and then went quiet",
           "There was no decline to detect. Entity_B's first-ever survey was already broken.",
           8, "Responders only. Entity_B joins at wave 2, Entity_C at wave 5, so neither has "
              "an earlier reading to decline from.", hl_size=25)
    picture(s, P2FIG / "14_purpose_trust_over_time.png",
            Inches(0.9), Inches(1.72), Inches(11.5), Inches(3.62))
    bullets(s, MARGIN, Inches(5.45), Inches(6.0), Inches(1.4), [
        "Entity_B's **first** measurement: purpose 2.975, trust 3.060, against "
        "NovaCorp-Origin's steady 3.38. A year later: 3.067. **Flat.**",
        "Survey response rate: Entity_B **62.6%** vs Entity_A **83.8%**, "
        "21 points lower.",
    ], size=12, gap=8)
    bullets(s, Inches(7.0), Inches(5.45), Inches(5.75), Inches(1.4), [
        "**307 employees were never issued a survey at all.** 152 of the mid-window "
        "leavers among them sit on the Entity_B and Entity_C legacy systems.",
        "They weren't silent. **They were never asked.**",
    ], size=12, gap=8)


def slide_09(prs):
    s = blank(prs)
    chrome(s, "You've already solved this once",
           "Entity_A is at 3.6%, below NovaCorp's own 4.5%. Integration works.", 9,
           "Attrition is annualised voluntary on active headcount. Response rate is the mean "
           "of each employee's own rate across five waves. Entity_A vs Entity_C, p = 0.101, "
           "not significant.")
    bullets(s, MARGIN, Inches(2.15), Inches(4.35), Inches(4.2), [
        "Entity_A (FY2022, fully integrated) is now the **safest cohort in the "
        "company**, better than NovaCorp-Origin itself.",
        "Its survey response rate has recovered to **83.8%**, indistinguishable "
        "from NovaCorp-Origin.",
        "Entity_C (late FY2024) is at **4.6%** and statistically indistinguishable "
        "from Entity_A. It is **not** a second Entity_B yet.",
    ], size=13, gap=13)
    # Purpose-built for this slide. A Kaplan-Meier curve was the obvious choice
    # and is wrong here: NovaCorp-Origin's curve is dominated by long-tenure
    # staff we only observe from 2024, so it reads far healthier than its
    # attrition rate, which contradicts this slide's headline. See A7.
    picture(s, DECKFIG / "S9_entity_recovery.png",
            Inches(5.35), Inches(1.95), Inches(7.35), Inches(4.50))
    callout(s, MARGIN, Inches(5.85), Inches(4.35), Inches(0.95),
            "This turns the recommendation from “we hope this helps” into "
            "“you have already done this once, in this company, with this "
            "workforce.” It also sets your success measure.", accent=TEAL, size=10.5)


def slide_10(prs):
    s = blank(prs)
    chrome(s, "What to do", "Three actions, in order. The first one is free.", 10)
    items = [
        ("1", "Redefine what counts as a regrettable loss", "~$0 · policy · this quarter",
         "Include High Performer alongside Outstanding. Separate the flag from the function "
         "that approved the exit. Review quarterly. Restates your baseline immediately and "
         "fixes next year's $42M at source.", PURPLE),
        ("2", "Finish the Entity_B integration, as a trust problem not a systems problem",
         "inside the guided $40–50M FY26 integration budget",
         "Migrate off BambooHR. Stabilise the Senior Manager layer first — Entity_B loses L3s "
         "at 9.8%/yr against Entity_A's 2.9%, a 3.4× gap — or there is nobody credible left to "
         "communicate the integration. Then leadership visibility and honest communication. "
         "Not manager training.", CORAL),
        ("3", "Measure acquired cohorts from day one, at cohort level", "~$0",
         "Baseline every acquired group against company norms at their first survey. Track "
         "response rate as a health metric in its own right. Fix the 307 who were never "
         "surveyed.", TEAL),
    ]
    y = Inches(2.00)
    for n, title, cost, body, accent in items:
        rect(s, MARGIN, y, Inches(0.055), Inches(1.32), fill=accent)
        tf = textbox(s, MARGIN + Inches(0.26), y - Inches(0.04), Inches(0.5), Inches(0.5))
        para(tf, n, 24, accent, bold=True, first=True, space_after=0)
        tf = textbox(s, MARGIN + Inches(0.80), y, Inches(7.6), Inches(0.35))
        para(tf, title, 15, DEEP, bold=True, first=True, space_after=0)
        tf = textbox(s, MARGIN + Inches(0.80), y + Inches(0.32), Inches(11.2), Inches(0.28))
        para(tf, cost, 10.5, accent, bold=True, italic=True, first=True, space_after=0)
        tf = textbox(s, MARGIN + Inches(0.80), y + Inches(0.62), Inches(11.3), Inches(0.7))
        para(tf, body, 12, INK, first=True, space_after=0, line_spacing=1.16)
        y += Inches(1.58)

    callout(s, MARGIN, Inches(6.72), Inches(12.0), Inches(0.5),
            "Notice what's not on this list: a pay round, a manager training programme, "
            "and cutting agency recruitment. The data rules out all three.", size=11.5)


def slide_11(prs):
    s = blank(prs)
    chrome(s, "What it's worth", "$71M addressable. The cheapest lever is the largest.", 11,
           "All figures from cost_model.py. The 20% and 40% reduction rates are assumptions, "
           "not findings. Entity_B's $6.2M is a subset of the $7.9M early-tenure lever, "
           "not additional to it.")
    rows = [
        ["Lever", "Addressable", "@20%", "@40%", "Cost to act"],
        ["Redefine regrettable_flag + quarterly review", "$45.1M", "$9.0M", "$18.1M", "~$0 (policy)"],
        ["Disengagement (813 staff below 2.5)", "$17.8M", "$3.6M", "$7.1M", "programme-dependent"],
        ["Early-tenure / acquisition onboarding", "$7.9M", "$1.6M", "$3.2M", "~$500–800/head"],
        ["      of which Entity_B specifically", "$6.2M", "$1.2M", "$2.5M", ""],
    ]
    table(s, MARGIN, Inches(2.05), Inches(12.1), rows,
          [Inches(4.55), Inches(1.75), Inches(1.55), Inches(1.55), Inches(2.70)],
          row_h=Inches(0.46), emphasis_rows=(1,))

    callout(s, MARGIN, Inches(4.35), Inches(12.0), Inches(0.95),
            "The 20% and 40% are assumptions, not findings. That's the band typically claimed "
            "for targeted retention programmes. We show both so the decision doesn't depend on "
            "the optimistic one. The first lever is a definition change, so its return isn't "
            "capped by budget at all.", size=11.5)

    rect(s, MARGIN, Inches(5.55), Inches(12.1), Inches(1.05), fill=RGBColor(0xFF, 0xF6, 0xE6))
    tf = textbox(s, MARGIN + Inches(0.3), Inches(5.72), Inches(11.5), Inches(0.75))
    para(tf, "$7.9M is the corrected figure.", 12, DEEP, bold=True, first=True, space_after=3)
    para(tf, "An earlier draft said $29.6M. That counted the baseline early attrition every "
             "employer carries, not the excess attributable to the acquisition cohorts. "
             "Method and restatement in appendix A7.", 11, INK, space_after=0,
         line_spacing=1.14)


def slide_12(prs):
    s = blank(prs)
    chrome(s, "How not to act on this",
           "A “who stopped answering” flag looks free. It flags 16.6% of 18–24s "
           "and 1.4% of 45–49s.", 12,
           "Four-fifths rule applied to a proposed individual flight-risk flag, on all active "
           "staff. All figures from ethics_audit.py. Full fairness method and results in "
           "appendix A4.", hl_size=25)
    bullets(s, MARGIN, Inches(2.05), Inches(4.9), Inches(4.5), [
        "The obvious move from our own analysis is a flight-risk score built on survey "
        "silence. **We recommend against deploying it at individual level.**",
        "Under the **four-fifths rule** it fails on age (ratio **0.08**, a twelvefold gap), "
        "on role level (L1 9.2% vs L5 1.4%) and on acquisition cohort (Entity_C 31.7% vs "
        "NovaCorp-Origin 2.4%).",
        "It is also **wrong most of the time**: **901 false positives** against 263 true "
        "ones, **precision 23%**. The dominant outcome is a career conversation about "
        "someone who was never leaving.",
        "Low response is concentrated in the acquisition cohorts, so the flag is "
        "substantially measuring **integration failure**, not individual intent.",
    ], size=12, gap=10)
    picture(s, EFIG / "E3_adverse_impact_age.png",
            Inches(5.85), Inches(1.95), Inches(6.85), Inches(4.60))
    callout(s, MARGIN, Inches(6.00), Inches(4.9), Inches(0.80),
            "We're showing you the thing we decided not to recommend. In a performance or "
            "redundancy conversation this is a discrimination exposure that costs more than "
            "the attrition does.", accent=CORAL, size=10.5)


def slide_13(prs):
    s = blank(prs)
    chrome(s, "How to use it responsibly",
           "Team-level diagnostic, disclosed to staff, never an individual score.", 13)
    items = [
        ("Never as an individual score handed to a manager.", "23% precision. Three in four "
         "flagged people were never leaving."),
        ("Team level only, minimum 8 responses.", "Reported as a diagnostic, not a prediction."),
        ("Tell employees the participation metric exists.", "If the survey was sold as "
         "confidential, quietly re-using non-response as a personal risk score breaks that "
         "promise. Once staff work it out, response rates collapse and you lose the "
         "instrument entirely."),
    ]
    y = Inches(2.15)
    for i, (t, b) in enumerate(items):
        tf = textbox(s, MARGIN, y, Inches(0.45), Inches(0.4))
        para(tf, f"{i+1}", 22, PURPLE, bold=True, first=True, space_after=0)
        tf = textbox(s, MARGIN + Inches(0.62), y + Inches(0.04), Inches(6.3), Inches(0.35))
        para(tf, t, 15, DEEP, bold=True, first=True, space_after=0)
        tf = textbox(s, MARGIN + Inches(0.62), y + Inches(0.42), Inches(6.3), Inches(0.85))
        para(tf, b, 12.5, INK, first=True, space_after=0, line_spacing=1.2)
        y += Inches(1.48)

    x = Inches(7.85)
    rect(s, x, Inches(2.10), Inches(4.85), Inches(4.15), fill=RGBColor(0xF7, 0xF3, 0xFC))
    tf = textbox(s, x + Inches(0.38), Inches(2.45), Inches(4.1), Inches(0.35))
    para(tf, "WHAT WE DELIBERATELY DID NOT DO", 10, GREY, bold=True, first=True,
         space_after=18)
    for t, b in [
        ("No demographic features in any model.",
         "Gender, age and cultural background appear only as fairness-audit dimensions, "
         "never as predictors."),
        ("No individual employees or managers named.",
         "Manager results aggregated, suppressed below 8 reports."),
        ("No causal language.", "Association only. We never claim one thing caused another."),
    ]:
        para(tf, t, 12.5, DEEP, bold=True, space_after=3)
        para(tf, b, 11, INK, space_after=17, line_spacing=1.18)


def slide_14(prs):
    s = blank(prs)
    chrome(s, "What we can't tell you",
           "The limits of this analysis, stated before you ask.", 14,
           "Full data-quality register in appendix A6. Statistical test register in A8.")
    left = [
        "**Two-year window, annualised.** All rates are voluntary exits ÷ 2 years ÷ active "
        "headcount = **4.7%/yr**. Your FY2025 Annual Report's “10.4% voluntary "
        "attrition” is total attrition including 267 involuntary exits across the full "
        "two years. It reproduces exactly as 1,400/13,403. **We restate; we don't accuse.**",
        "**We cannot see pre-acquisition Entity_B.** They may have arrived unhappy from their "
        "previous employer. We genuinely cannot rule this out.",
        "**Engagement is only observed for responders**, and leavers responded less, so every "
        "disengagement figure is a lower bound.",
        "**Association, not causation.** Senior Manager churn and collapsed trust move "
        "together. We cannot say which drives which.",
    ]
    right = [
        "**regrettable_flag, performance_band_at_exit and stated_exit_reason are retrospective "
        "HR judgements.** We use them as the object of analysis, not as ground truth.",
        "**The 2025-H2 review cycle is missing**, so performance data is up to 12 months stale "
        "for some leavers.",
        "**Small samples above Level 4.** Nothing above L4 is reliable: 279 people at L4, "
        "103 above it.",
    ]
    bullets(s, MARGIN, Inches(1.95), Inches(6.0), Inches(4.5), left, size=11.5, gap=13)
    bullets(s, Inches(7.05), Inches(1.95), Inches(5.7), Inches(4.5), right, size=11.5, gap=13)

    callout(s, MARGIN, Inches(6.05), Inches(12.0), Inches(0.62),
            "On a strict voluntary reading you have already met the FY2026 sub-9.5% target. "
            "That's worth resolving before the next board pack.", accent=GOLD, size=11.5)


def slide_15(prs):
    s = blank(prs)
    chrome(s, "Close", "Start with the free one.", 15)
    items = [
        ("THIS QUARTER", "Redefine regrettable.", PURPLE,
         "Cost ~$0. It restates your baseline and stops you funding a budget line built on a "
         "blind spot."),
        ("THIS HALF", "Finish Entity_B.", CORAL,
         "Senior Manager layer first, then leadership visibility. $8.4M concentrated, inside a "
         "budget you have already guided."),
        ("ONGOING", "Baseline every acquired cohort from day one.", TEAL,
         "Entity_C is healthy today. It does not have to become the next Entity_B."),
    ]
    x = MARGIN
    cw = Inches(3.87)
    for label, title, accent, body in items:
        rect(s, x, Inches(2.05), cw, Inches(0.06), fill=accent)
        tf = textbox(s, x, Inches(2.32), cw, Inches(0.3))
        para(tf, label, 10.5, accent, bold=True, first=True, space_after=10)
        tf = textbox(s, x, Inches(2.68), cw, Inches(0.8))
        para(tf, title, 19, DEEP, bold=True, first=True, space_after=0, line_spacing=1.06)
        tf = textbox(s, x, Inches(3.72), cw, Inches(1.3))
        para(tf, body, 12.5, INK, first=True, space_after=0, line_spacing=1.2)
        x += cw + Inches(0.24)

    rect(s, MARGIN, Inches(5.35), Inches(12.1), Inches(1.35), fill=DEEP)
    tf = textbox(s, MARGIN + Inches(0.45), Inches(5.62), Inches(11.2), Inches(0.85))
    para(tf, "You already have the money set aside.", 15, WHITE, bold=True, first=True,
         space_after=5)
    para(tf, "The Annual Report commits to a People Reinvention Programme three times. "
             "This is what we'd spend it on, and this is what we'd measure.", 12.5,
         RGBColor(0xD8, 0xC2, 0xEF), space_after=0, line_spacing=1.16)


# ==========================================================================
def main():
    prs = new_deck()
    for fn in [slide_01, slide_02, slide_03, slide_04, slide_05, slide_06,
               slide_07, slide_08, slide_09, slide_10, slide_11, slide_12,
               slide_13, slide_14, slide_15]:
        fn(prs)

    assert len(prs.slides) == 15, f"main deck must be 15 slides, got {len(prs.slides)}"
    out = OUT_DIR / "NovaCorp_deck.pptx"
    prs.save(out)
    print(f"  {len(prs.slides)} slides -> {out}")
    if TEAM_NAME.startswith("TEAM NAME"):
        print("  WARNING: team name is still the placeholder, set TEAM_NAME before export")


if __name__ == "__main__":
    main()
